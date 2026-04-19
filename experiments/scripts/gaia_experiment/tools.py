"""
GAIA experiment tools.

Tools provided:
  google_web_search  - web search via SerpAPI / Serper (reused from finagent_experiment)
  calculator         - safe arithmetic expression evaluator
  file_reader        - reads PDF / Excel / CSV / Word / PowerPoint / plain-text files
  python_executor    - executes Python code in an isolated subprocess
  multimodal_viewer  - analyses images / audio / video via the Doubao multimodal API
"""

import base64
import io
import logging
import math
import os
import re
import subprocess
import sys
import tempfile
from typing import Any, Dict, List

from finagent_experiment.tools import GoogleWebSearch, FinancialTool

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Calculator
# ---------------------------------------------------------------------------

class Calculator(FinancialTool):
    """Safe arithmetic expression evaluator backed by Python's math module."""

    name: str = "calculator"
    description: str = (
        "Evaluate a mathematical expression and return the numeric result. "
        "Supports standard arithmetic (+, -, *, /, **), parentheses, and math functions: "
        "sqrt, log, log10, sin, cos, tan, abs, round, ceil, floor, pi, e. "
        "Examples: '(3 + 4) * 2', 'sqrt(144)', 'log(100, 10)'."
    )
    input_arguments: Dict[str, Any] = {
        "expression": {
            "type": "string",
            "description": "The mathematical expression to evaluate.",
        }
    }
    required_arguments: List[str] = ["expression"]

    _SAFE_GLOBALS: Dict[str, Any] = {
        "__builtins__": {},
        "abs": abs, "round": round, "min": min, "max": max, "sum": sum,
        "sqrt": math.sqrt, "log": math.log, "log10": math.log10, "log2": math.log2,
        "exp": math.exp, "pow": math.pow,
        "sin": math.sin, "cos": math.cos, "tan": math.tan,
        "asin": math.asin, "acos": math.acos, "atan": math.atan, "atan2": math.atan2,
        "ceil": math.ceil, "floor": math.floor,
        "pi": math.pi, "e": math.e, "inf": math.inf,
    }

    async def call_tool(self, arguments: dict, **kwargs) -> str:
        expression = arguments.get("expression", "").strip()
        if not expression:
            raise ValueError("No expression provided.")
        if not re.match(r"^[0-9\s\+\-\*\/\(\)\.\,\^\_a-zA-Z]+$", expression):
            raise ValueError(f"Expression contains disallowed characters: {expression!r}")
        result = eval(expression.replace("^", "**"), self._SAFE_GLOBALS, {})  # noqa: S307
        return str(result)


# ---------------------------------------------------------------------------
# FileReader  (PDF / Excel / CSV / Word / PowerPoint / plain text)
# ---------------------------------------------------------------------------

class FileReader(FinancialTool):
    """
    Read and extract text content from local files.

    Supported formats:
      .pdf              — via pdfplumber
      .xlsx / .xls      — via openpyxl (returns sheet data as markdown tables)
      .csv              — via csv module
      .docx             — via python-docx
      .pptx             — via python-pptx
      .txt / .md / .py / .json / .jsonld — plain UTF-8 read

    The file path must be an absolute local path or a path relative to the
    working directory. Use the local_file_path field injected into each sample
    by the runner, or construct the path from sample_info.file_path.
    """

    name: str = "file_reader"
    description: str = (
        "Read and extract text content from a local file. "
        "Supported formats: PDF (.pdf), Excel (.xlsx/.xls), CSV (.csv), "
        "Word (.docx), PowerPoint (.pptx), and plain text (.txt, .py, .json, etc.). "
        "Returns the extracted text content."
    )
    input_arguments: Dict[str, Any] = {
        "file_path": {
            "type": "string",
            "description": "Absolute or relative path to the local file to read.",
        },
        "max_chars": {
            "type": "integer",
            "description": "Maximum number of characters to return (default 20000).",
        },
    }
    required_arguments: List[str] = ["file_path"]

    async def call_tool(self, arguments: dict, **kwargs) -> str:
        file_path = arguments.get("file_path", "").strip()
        max_chars = int(arguments.get("max_chars", 20000))

        if not file_path:
            raise ValueError("file_path is required.")
        if file_path.startswith("http://") or file_path.startswith("https://"):
            return (
                "ERROR: file_reader requires a local file path, not a URL. "
                "The local file path for this sample is provided in the question as 'Local Path'. "
                "For remote URLs, use google_web_search or multimodal_viewer instead."
            )
        if not os.path.exists(file_path):
            file_path = _resolve_path_case(file_path)
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        ext = os.path.splitext(file_path)[1].lower()
        text = self._read(file_path, ext)

        if len(text) > max_chars:
            text = text[:max_chars] + f"\n... [truncated at {max_chars} chars, total {len(text)}]"
        return text

    def _read(self, path: str, ext: str) -> str:
        if ext == ".pdf":
            return self._read_pdf(path)
        elif ext in (".xlsx", ".xls"):
            return self._read_excel(path)
        elif ext == ".csv":
            return self._read_csv(path)
        elif ext == ".docx":
            return self._read_docx(path)
        elif ext == ".pptx":
            return self._read_pptx(path)
        else:
            # Plain text fallback (txt, py, md, json, jsonld, pdb, etc.)
            return self._read_text(path)

    def _read_pdf(self, path: str) -> str:
        try:
            import pdfplumber
        except ImportError:
            raise ImportError("pdfplumber is required for PDF reading: pip install pdfplumber")
        pages = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"--- Page {i + 1} ---\n{text}")
        return "\n\n".join(pages) if pages else "(No extractable text found in PDF)"

    def _read_excel(self, path: str) -> str:
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required for Excel reading: pip install openpyxl")
        wb = openpyxl.load_workbook(path, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                # Skip entirely empty rows
                if any(cell is not None for cell in row):
                    rows.append(" | ".join("" if c is None else str(c) for c in row))
            if rows:
                sheets.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
        return "\n\n".join(sheets) if sheets else "(Empty workbook)"

    def _read_csv(self, path: str) -> str:
        import csv
        rows = []
        with open(path, newline="", encoding="utf-8-sig") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" | ".join(row))
        return "\n".join(rows) if rows else "(Empty CSV)"

    def _read_docx(self, path: str) -> str:
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required: pip install python-docx")
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _read_pptx(self, path: str) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required: pip install python-pptx")
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
        return "\n\n".join(slides) if slides else "(No text found in presentation)"

    def _read_text(self, path: str) -> str:
        for enc in ("utf-8", "utf-8-sig", "gbk", "latin-1"):
            try:
                with open(path, encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise ValueError(f"Could not decode file with any known encoding: {path}")


# ---------------------------------------------------------------------------
# PythonExecutor
# ---------------------------------------------------------------------------

class PythonExecutor(FinancialTool):
    """
    Execute Python code in an isolated subprocess and return stdout + stderr.

    The code runs with a configurable timeout (default 30 s) in a fresh
    subprocess using the same Python interpreter as the host.  No network or
    filesystem sandboxing is applied — use with trusted LLM-generated code only.
    """

    name: str = "python_executor"
    description: str = (
        "Execute Python code and return its stdout output. "
        "Use this for data analysis, calculations, file processing, or any task "
        "that benefits from running Python. "
        "The code has access to standard library modules. "
        "Print your result to stdout so it is captured. "
        "Example: print(sum(range(1, 101)))"
    )
    input_arguments: Dict[str, Any] = {
        "code": {
            "type": "string",
            "description": "The Python code to execute. Print results to stdout.",
        },
        "timeout": {
            "type": "integer",
            "description": "Execution timeout in seconds (default 30, max 120).",
        },
    }
    required_arguments: List[str] = ["code"]

    async def call_tool(self, arguments: dict, **kwargs) -> str:
        code = arguments.get("code", "")
        timeout = min(int(arguments.get("timeout", 30)), 120)

        if not code.strip():
            raise ValueError("No code provided.")

        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".py", delete=False, encoding="utf-8"
        ) as tmp:
            tmp.write(code)
            tmp_path = tmp.name

        try:
            proc = subprocess.run(
                [sys.executable, tmp_path],
                capture_output=True,
                text=True,
                timeout=timeout,
            )
            stdout = proc.stdout.strip()
            stderr = proc.stderr.strip()
            if stderr:
                return f"{stdout}\n[stderr]\n{stderr}" if stdout else f"[stderr]\n{stderr}"
            return stdout if stdout else "(No output)"
        except subprocess.TimeoutExpired:
            raise TimeoutError(f"Code execution timed out after {timeout}s.")
        finally:
            os.unlink(tmp_path)


# ---------------------------------------------------------------------------
# MultimodalViewer  (image / audio / video via Doubao API)
# ---------------------------------------------------------------------------

class MultimodalViewer(FinancialTool):
    """
    Analyse images, audio, or video files using the Doubao multimodal API.

    Calls the Doubao model (doubao-seed-2-0-lite-260215 by default) via the
    OpenAI-compatible client at https://ark.cn-beijing.volces.com/api/v3.

    For local files the content is base64-encoded and passed as a data URI.
    For remote URLs the URL is passed directly.

    Supported local file types:
      Images : .png, .jpg, .jpeg, .gif, .webp, .bmp
      Audio  : .mp3, .wav, .m4a, .ogg, .flac
      Video  : .mp4, .mov, .avi, .mkv, .webm
    """

    name: str = "multimodal_viewer"
    description: str = (
        "Analyse the content of an image, audio, or video file using a multimodal AI model. "
        "Provide a local file path or a public URL. "
        "Use this to describe images, transcribe audio, or summarise video content. "
        "Specify what you want to know about the file in the 'prompt' argument."
    )
    input_arguments: Dict[str, Any] = {
        "file_path": {
            "type": "string",
            "description": (
                "Local file path or public URL of the image/audio/video to analyse. "
                "For local files, use the local_file_path field from the sample."
            ),
        },
        "prompt": {
            "type": "string",
            "description": "What you want to know about the file (e.g. 'Describe this image in detail').",
        },
    }
    required_arguments: List[str] = ["file_path", "prompt"]

    # MIME type mapping
    _MIME: Dict[str, str] = {
        ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp",
        ".mp3": "audio/mpeg", ".wav": "audio/wav", ".m4a": "audio/mp4",
        ".ogg": "audio/ogg", ".flac": "audio/flac",
        ".mp4": "video/mp4", ".mov": "video/quicktime", ".avi": "video/x-msvideo",
        ".mkv": "video/x-matroska", ".webm": "video/webm",
    }

    _IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
    _AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac"}
    _VIDEO_EXTS = {".mp4", ".mov", ".avi", ".mkv", ".webm"}

    def __init__(
        self,
        api_key: str = None,
        model: str = "doubao-seed-2-0-lite-260215",
        base_url: str = "https://ark.cn-beijing.volces.com/api/v3",
    ):
        self.api_key = api_key or os.getenv("ARK_API_KEY")
        self.model = model
        self.base_url = base_url

    def _get_client(self):
        try:
            from openai import OpenAI
        except ImportError:
            raise ImportError("openai package is required: pip install openai")
        if not self.api_key:
            raise ValueError("ARK_API_KEY is not set. Doubao multimodal API requires it.")
        from openai import OpenAI
        return OpenAI(base_url=self.base_url, api_key=self.api_key)

    def _is_url(self, path: str) -> bool:
        return path.startswith("http://") or path.startswith("https://")

    def _to_data_uri(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        mime = self._MIME.get(ext, "application/octet-stream")
        with open(file_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
        return f"data:{mime};base64,{b64}"

    def _detect_media_type(self, path: str) -> str:
        ext = os.path.splitext(path)[1].lower()
        if ext in self._IMAGE_EXTS:
            return "image"
        if ext in self._AUDIO_EXTS:
            return "audio"
        if ext in self._VIDEO_EXTS:
            return "video"
        return "image"  # default fallback

    async def call_tool(self, arguments: dict, **kwargs) -> str:
        file_path = arguments.get("file_path", "").strip()
        prompt = arguments.get("prompt", "Describe the content of this file in detail.")

        if not file_path:
            raise ValueError("file_path is required.")

        is_url = self._is_url(file_path)
        if not is_url:
            if not os.path.exists(file_path):
                file_path = _resolve_path_case(file_path)
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")

        media_type = self._detect_media_type(file_path)
        image_url = file_path if is_url else self._to_data_uri(file_path)

        client = self._get_client()

        # Build content blocks according to media type
        if media_type == "image":
            content = [
                {"type": "input_image", "image_url": image_url},
                {"type": "input_text", "text": prompt},
            ]
        elif media_type == "audio":
            content = [
                {"type": "input_audio", "audio_url": image_url},
                {"type": "input_text", "text": prompt},
            ]
        else:  # video
            content = [
                {"type": "input_video", "video_url": image_url},
                {"type": "input_text", "text": prompt},
            ]

        response = client.responses.create(
            model=self.model,
            input=[{"role": "user", "content": content}],
        )

        # Extract text from response
        if hasattr(response, "output_text"):
            return response.output_text
        if hasattr(response, "output") and response.output:
            for item in response.output:
                if hasattr(item, "content"):
                    for c in item.content:
                        if hasattr(c, "text"):
                            return c.text
        return str(response)


# ---------------------------------------------------------------------------
# Factory
# ---------------------------------------------------------------------------

def _resolve_path_case(file_path: str) -> str:
    """Try to fix case-sensitivity issues in file paths (e.g. Validation vs validation)."""
    if os.path.exists(file_path):
        return file_path
    parts = file_path.replace("\\", "/").split("/")
    resolved = parts[0] if parts[0] else "/"
    for part in parts[1:]:
        parent = resolved if resolved else "."
        try:
            entries = os.listdir(parent)
        except OSError:
            return file_path
        match = next((e for e in entries if e.lower() == part.lower()), None)
        resolved = os.path.join(resolved, match if match else part)
    return resolved


def get_all_tool_definitions(tools: Dict[str, "FinancialTool"]) -> List[Dict[str, Any]]:
    """Build OpenAI-format function schema list from tool instances."""
    definitions = []
    for tool in tools.values():
        tool_def = tool.get_tool_definition()
        definitions.append({
            "type": "function",
            "function": {
                "name": tool_def.body["name"],
                "description": tool_def.body["description"],
                "parameters": {
                    "type": "object",
                    "properties": tool_def.body["properties"],
                    "required": tool_def.body["required"],
                },
            },
        })
    return definitions


def create_gaia_tools(
    serpapi_api_key: str = None,
    serper_api_key: str = None,
    ark_api_key: str = None,
    doubao_model: str = "doubao-seed-2-0-lite-260215",
) -> Dict[str, FinancialTool]:
    """
    Create the full GAIA tool set.

    Args:
        serpapi_api_key: SerpAPI key (primary web search); falls back to SERPAPI_API_KEY env var
        serper_api_key:  Serper API key (fallback web search); falls back to SERPER_API_KEY env var
        ark_api_key:     Doubao / ByteDance ARK API key; falls back to ARK_API_KEY env var
        doubao_model:    Doubao model ID for multimodal analysis

    Returns:
        Dict mapping tool name -> tool instance
    """
    return {
        "google_web_search": GoogleWebSearch(
            serpapi_api_key=serpapi_api_key or os.getenv("SERPAPI_API_KEY"),
            serper_api_key=serper_api_key or os.getenv("SERPER_API_KEY"),
        ),
        "calculator": Calculator(),
        "file_reader": FileReader(),
        "python_executor": PythonExecutor(),
        "multimodal_viewer": MultimodalViewer(
            api_key=ark_api_key or os.getenv("ARK_API_KEY"),
            model=doubao_model,
        ),
    }
